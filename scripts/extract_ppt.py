#!/usr/bin/env python3
"""
提取 PPT (.pptx) 文件中每页的文字内容和备注。
用法: python extract_ppt.py <input.pptx> [output.json]
"""
import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("请先安装 python-pptx: pip install python-pptx")
    sys.exit(1)


def extract_ppt(ppt_path: str) -> list[dict]:
    prs = Presentation(ppt_path)
    slides = []

    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        notes_text = ""

        # 提取正文文字
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

        # 提取备注
        if slide.has_notes_slide:
            notes = slide.notes_slide
            notes_text = notes.notes_text_frame.text.strip()

        slides.append({
            "slide": idx,
            "content": "\n".join(texts),
            "notes": notes_text
        })

    return slides


def main():
    if len(sys.argv) < 2:
        print("用法: python extract_ppt.py <input.pptx> [output.json]")
        sys.exit(1)

    ppt_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    if not Path(ppt_path).exists():
        print(f"错误: 文件不存在 — {ppt_path}")
        sys.exit(1)

    slides = extract_ppt(ppt_path)

    # 生成摘要
    result = {
        "total_slides": len(slides),
        "slides": slides,
        "full_text": "\n\n--- 幻灯片分隔 ---\n\n".join(
            f"[幻灯片 {s['slide']}]\n{s['content']}\n[备注]: {s['notes']}"
            for s in slides
        )
    }

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"已输出到: {output_path}")
    else:
        # 直接输出全文
        print(f"=== 共 {len(slides)} 页幻灯片 ===\n")
        print(result["full_text"])


if __name__ == "__main__":
    main()
